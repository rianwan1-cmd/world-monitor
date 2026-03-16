from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

print(">>> 正在生成面向大专生的高能活泼版《公关实务》课件...")

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
blank_layout = prs.slide_layouts[6]

# 颜色定义：深邃黑底 + 赛博霓虹色点缀（大专生/Z世代审美）
BG_BLACK = RGBColor(18, 18, 20)
TEXT_WHITE = RGBColor(245, 245, 245)
TEXT_GRAY = RGBColor(150, 150, 160)
NEON_GREEN = RGBColor(57, 255, 20)
NEON_PINK = RGBColor(255, 42, 127)
NEON_BLUE = RGBColor(0, 243, 255)
WARNING_YELLOW = RGBColor(255, 204, 0)

def set_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = 'Microsoft YaHei' # 使用微软雅黑保证兼容和现代感
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return p

# ==========================================
# 幻灯片 1：高能封面
# ==========================================
slide1 = prs.slides.add_slide(blank_layout)
set_background(slide1, BG_BLACK)
# 装饰线
line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(16), Inches(0.2))
line.fill.solid()
line.fill.fore_color.rgb = NEON_GREEN
line.line.fill.background()

add_text(slide1, Inches(1), Inches(3), Inches(14), Inches(2), "🔥 公共关系实务", 100, TEXT_WHITE, True, PP_ALIGN.CENTER)
add_text(slide1, Inches(1), Inches(5), Inches(14), Inches(1), "不仅教你怎么说话，更教你怎么“搞事情”", 36, NEON_PINK, False, PP_ALIGN.CENTER)
add_text(slide1, Inches(1), Inches(6.5), Inches(14), Inches(1), "教材：霍瑞红 (第四版)   |   核心 17 课时高能导学", 24, TEXT_GRAY, False, PP_ALIGN.CENTER)

# ==========================================
# 幻灯片 2：17课时 升级打怪地图
# ==========================================
slide2 = prs.slides.add_slide(blank_layout)
set_background(slide2, BG_BLACK)
add_text(slide2, Inches(1), Inches(0.8), Inches(14), Inches(1), "🗺️ 这学期，我们怎么“打怪升级”？", 54, NEON_BLUE, True)

# 分三列展示 17 个课时核心模块
col_width = Inches(4.5)
# 列 1
add_text(slide2, Inches(1), Inches(2.2), col_width, Inches(0.5), "🟢 第一阶段：看透套路", 30, NEON_GREEN, True)
texts1 = ["课时 1: 公关是玄学还是科学？", "课时 2: 谁在忽悠？忽悠谁？", "课时 3: 寻找你的“上帝”与黑粉", "课时 4: 把话说到心坎里", "课时 5: 双微一抖的新媒体战"]
for i, t in enumerate(texts1):
    add_text(slide2, Inches(1), Inches(3 + i*0.8), col_width, Inches(0.5), t, 22, TEXT_WHITE)

# 列 2
add_text(slide2, Inches(6), Inches(2.2), col_width, Inches(0.5), "🟡 第二阶段：动手搞事", 30, WARNING_YELLOW, True)
texts2 = ["课时 6: 无调查，不公关", "课时 7: 写个不被毙的策划案", "课时 8: 活动落地与避坑指南", "课时 9: 老板的钱花哪了？", "课时 10: 爆款发布会是怎么炼成的", "课时 11: 联名：花钱听个响？"]
for i, t in enumerate(texts2):
    add_text(slide2, Inches(6), Inches(3 + i*0.8), col_width, Inches(0.5), t, 22, TEXT_WHITE)

# 列 3
add_text(slide2, Inches(11), Inches(2.2), col_width, Inches(0.5), "🔴 第三阶段：绝地求生", 30, NEON_PINK, True)
texts3 = ["课时 12: 老板上热搜了怎么办！", "课时 13: 满分道歉信怎么写？", "课时 14: 键盘侠退散 (模拟战)", "课时 15: 公关人的神仙职场素养", "课时 16: 舌战群儒与宴请礼仪", "课时 17: 期末大考：真题提案"]
for i, t in enumerate(texts3):
    add_text(slide2, Inches(11), Inches(3 + i*0.8), col_width, Inches(0.5), t, 22, TEXT_WHITE)

# ==========================================
# 幻灯片 3：课时 1 破冰 - 灵魂拷问
# ==========================================
slide3 = prs.slides.add_slide(blank_layout)
set_background(slide3, BG_BLACK)
add_text(slide3, Inches(1), Inches(1), Inches(14), Inches(1.5), "🤔 课时 01：灵魂拷问", 64, WARNING_YELLOW, True, PP_ALIGN.CENTER)
add_text(slide3, Inches(1), Inches(3), Inches(14), Inches(1.5), "在你眼里，什么是“公关（PR）”？", 50, TEXT_WHITE, True, PP_ALIGN.CENTER)

add_text(slide3, Inches(2), Inches(5), Inches(5), Inches(2), "❌ 刻板印象\n\n喝酒应酬、拉关系\n遇到丑闻发通稿\n花钱雇水军删帖", 32, TEXT_GRAY, False, PP_ALIGN.CENTER)
add_text(slide3, Inches(9), Inches(5), Inches(5), Inches(2), "✅ 真实公关\n\n管理品牌认知\n建立信任关系\n转危为机的战略家", 32, NEON_GREEN, True, PP_ALIGN.CENTER)

# ==========================================
# 幻灯片 4：案例引入（瑞幸）
# ==========================================
slide4 = prs.slides.add_slide(blank_layout)
set_background(slide4, BG_BLACK)
add_text(slide4, Inches(1), Inches(1), Inches(14), Inches(1), "☕ 真实商战吃瓜：瑞幸起死回生", 54, NEON_BLUE, True)

add_text(slide4, Inches(1), Inches(2.5), Inches(14), Inches(1), "2020年，瑞幸因财务造假被美股退市，面临全网声讨，陷入绝境。它怎么活下来的？", 32, TEXT_WHITE, False)

add_text(slide4, Inches(1), Inches(4.2), Inches(14), Inches(4), "👉 认错干脆：第一时间承认，态度诚恳（危机公关 5S 原则）。\n\n👉 制造爆款：推出“生椰拿铁”，用硬产品转移注意力（策划与传播）。\n\n👉 自黑玩梗：在社交媒体上化身“打工人”，拉近与年轻人的距离（沟通与客体分析）。\n\n👉 跨界联名：与茅台联名“酱香拿铁”，彻底引爆全网（专题活动策划）。", 28, TEXT_GRAY, False)

# ==========================================
# 幻灯片 5：课程考核规则
# ==========================================
slide5 = prs.slides.add_slide(blank_layout)
set_background(slide5, BG_BLACK)
add_text(slide5, Inches(1), Inches(1), Inches(14), Inches(1.5), "🎮 游戏规则（考核标准）", 64, NEON_PINK, True, PP_ALIGN.CENTER)

add_text(slide5, Inches(2), Inches(3.5), Inches(12), Inches(3), "无聊的死记硬背？不存在的！\n这门课，我们用【实战打分】", 40, TEXT_WHITE, True, PP_ALIGN.CENTER)

add_text(slide5, Inches(1), Inches(6), Inches(4.5), Inches(2), "🗣️ 课堂发言 & 提案\n30%", 32, NEON_BLUE, True, PP_ALIGN.CENTER)
add_text(slide5, Inches(5.75), Inches(6), Inches(4.5), Inches(2), "📝 分组实战项目\n30%", 32, WARNING_YELLOW, True, PP_ALIGN.CENTER)
add_text(slide5, Inches(10.5), Inches(6), Inches(4.5), Inches(2), "🏆 期末路演大考\n40%", 32, NEON_GREEN, True, PP_ALIGN.CENTER)

# ==========================================
# 幻灯片 6：课后任务
# ==========================================
slide6 = prs.slides.add_slide(blank_layout)
set_background(slide6, BG_BLACK)
add_text(slide6, Inches(1), Inches(2), Inches(14), Inches(1.5), "🎯 随堂新手村任务", 64, WARNING_YELLOW, True, PP_ALIGN.CENTER)
add_text(slide6, Inches(2), Inches(4), Inches(12), Inches(3), "1. 自由组队（每组 4-5 人），选出一个公关总监 (组长)。\n\n2. 寻找一个最近一年内发生的热门品牌事件（翻车或爆火均可）。\n\n3. 下节课，派代表上台分享：事件中的“主角”是谁？“受众”是谁？他们用了什么“传播手段”？", 32, TEXT_WHITE, False, PP_ALIGN.LEFT)

output_file = 'Huo_PR_17_Lessons_Vocational.pptx'
prs.save(output_file)
print(f">>> 高能版课件生成成功：{output_file}")
