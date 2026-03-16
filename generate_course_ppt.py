from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

print(">>> 启动 Apple Masterclass 课件渲染引擎...")

# 1. 渲染光追级渐变背景
def render_fluid_bg(filename, base_color, orbs):
    width, height = 1920, 1080
    bg = Image.new('RGB', (width, height), base_color)
    glow = Image.new('RGBA', (width, height), (0, 0, 0, 0))
    draw = ImageDraw.Draw(glow)
    for orb in orbs:
        x, y, r, color = orb
        draw.ellipse((x-r, y-r, x+r, y+r), fill=color)
    glow = glow.filter(ImageFilter.GaussianBlur(300))
    bg.paste(glow, (0,0), glow)
    bg.save(filename)

# 课件专属情绪配色
bgs = {
    "cover": ((5, 5, 12), [(960, 540, 800, (88, 86, 214, 150)), (400, 200, 800, (90, 200, 250, 100))]),
    "agenda": ((8, 8, 8), [(1600, 100, 700, (100, 100, 120, 80))]),
    "ch1": ((0, 10, 15), [(-200, 540, 900, (50, 150, 200, 120))]), # 沉稳蓝(定义)
    "ch2": ((10, 5, 0), [(1600, 540, 1000, (255, 149, 0, 120))]),  # 价值金(品牌)
    "ch3": ((0, 15, 5), [(1600, 540, 1000, (52, 199, 89, 100))]),  # 生态绿(媒介)
    "ch4": ((15, 0, 0), [(1600, 540, 1000, (255, 59, 48, 140))]),  # 警示红(危机)
    "quote": ((0, 0, 0), [(960, 1080, 800, (255, 255, 255, 40))])
}

print("-> 正在生成教学版流光背景...")
for name, data in bgs.items():
    render_fluid_bg(f"bg_{name}.png", data[0], data[1])

# 2. 构建 PPTX
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
blank_layout = prs.slide_layouts[6]

WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(160, 160, 165)
DARK_GRAY = RGBColor(40, 40, 45)
GOLD = RGBColor(212, 175, 55)

def add_text(slide, left, top, width, height, text, font_size, color, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = 'Helvetica Neue'
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return p

# --- 封面页 ---
slide = prs.slides.add_slide(blank_layout)
slide.shapes.add_picture("bg_cover.png", 0, 0, width=Inches(16), height=Inches(9))
add_text(slide, Inches(1), Inches(3.2), Inches(14), Inches(2), "公共关系实务", 90, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.2), Inches(14), Inches(1), "Public Relations Practice", 36, GOLD, False, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(6.5), Inches(14), Inches(1), "卓越品牌背后的战略基石 · 核心大师课", 24, GRAY, False, PP_ALIGN.CENTER)

# --- 导航页 (Agenda) ---
slide = prs.slides.add_slide(blank_layout)
slide.shapes.add_picture("bg_agenda.png", 0, 0, width=Inches(16), height=Inches(9))
add_text(slide, Inches(2), Inches(1.5), Inches(12), Inches(1), "课程核心模块", 54, WHITE, True)
add_text(slide, Inches(2), Inches(2.5), Inches(12), Inches(0.5), "Course Agenda", 24, GOLD, False)

agenda_items = [
    "01 / 认知重塑：公关的本质与核心边界",
    "02 / 品牌心智：形象塑造与声誉长期主义",
    "03 / 媒介生态：PESO 模型与全链路沟通",
    "04 / 风暴防御：危机公关与 5S 黄金原则"
]
for i, item in enumerate(agenda_items):
    add_text(slide, Inches(2.5), Inches(4 + i*1.2), Inches(10), Inches(1), item, 32, GRAY, False)

# --- 课件内容页 (带专业术语的巨型数字版) ---
def add_course_slide(bg_name, number, title, subtitle, points):
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.add_picture(f"bg_{bg_name}.png", 0, 0, width=Inches(16), height=Inches(9))
    
    # 背景巨型浮水印数字
    add_text(slide, Inches(0.5), Inches(0.5), Inches(8), Inches(8), number, 400, DARK_GRAY, True)
    
    # 标题区
    add_text(slide, Inches(6.5), Inches(1.5), Inches(8.5), Inches(1.5), title, 60, WHITE, True)
    add_text(slide, Inches(6.5), Inches(2.8), Inches(8.5), Inches(1), subtitle, 32, GOLD, False)
    
    # 知识点区 (使用多段落)
    txBox = slide.shapes.add_textbox(Inches(6.5), Inches(4.5), Inches(8.5), Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, pt in enumerate(points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = pt
        p.font.name = 'Helvetica Neue'
        p.font.size = Pt(26)
        p.font.color.rgb = WHITE if ":" in pt or "：" in pt else GRAY
        p.space_after = Pt(20) # 段间距增加呼吸感
        p.line_spacing = 1.3

# 第一章
add_course_slide("ch1", "01", "公关的本质", "管理公众认知，而非单纯制造信息", [
    "• 核心定义：PR 是组织与公众之间建立互信的战略沟通过程。",
    "• 认知误区：公关不是发通稿，也不是掩饰错误的“灭火器”。",
    "• 核心目标：让受众感受到“共鸣”，获取第三方背书，建立长期的信任关系（Trust & Relationship）。"
])

# 第二章
add_course_slide("ch2", "02", "形象与声誉", "品牌是你不在场时，别人的评价", [
    "• 知名度 vs 美誉度：知名度决定品牌能走多快，美誉度决定品牌能走多远。",
    "• 无形资产：声誉是企业在危机中最厚实的护城河。",
    "• PR Campaign：通过持续的战略性公关战役，输出一致的价值观，积累品牌资产。"
])

# 第三章
add_course_slide("ch3", "03", "媒介生态学", "理解 PESO 传播模型", [
    "• Paid Media (付费媒体)：广告、KOL投放，保证触达率。",
    "• Earned Media (赢得媒体)：新闻报道、口碑推荐，公关的核心，信任度最高。",
    "• Shared Media (共享媒体)：社交媒体裂变、用户UGC。",
    "• Owned Media (自有媒体)：官网、官方双微，控制权最强。"
])

# 第四章
add_course_slide("ch4", "04", "危机公关", "风暴中心的 5S 黄金原则", [
    "危机是品牌的放大器，也是试金石。处理危机的核心逻辑：",
    "1. 承担责任 (Shoulder)：态度重于事实，不推卸责任。",
    "2. 真诚沟通 (Sincerity)：拒绝官方套话，用人情味化解对立。",
    "3. 速度第一 (Speed)：打破黄金24小时，争取“黄金1小时”响应。",
    "4. 系统运行 (System)：统一对外口径，内部协同作战。",
    "5. 权威证实 (Standard)：引入第三方权威机构背书。"
])

# --- 总结页 ---
slide = prs.slides.add_slide(blank_layout)
slide.shapes.add_picture("bg_quote.png", 0, 0, width=Inches(16), height=Inches(9))
add_text(slide, Inches(2), Inches(3.5), Inches(12), Inches(2), "“公关的最高境界，\n是让公众成为你的坚定同盟。”", 50, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(2), Inches(6), Inches(12), Inches(1), "— 课程总结 —", 24, GRAY, False, PP_ALIGN.CENTER)

output_file = 'PR_Practice_Courseware.pptx'
prs.save(output_file)
print(f">>> 教学版课件渲染完成！文件已保存为: {output_file}")
