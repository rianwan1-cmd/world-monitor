from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

print("1. 正在渲染流光溢彩的高级光晕背景图...")
# 设置画布尺寸 1920x1080 (16:9)
width, height = 1920, 1080
# 极暗的深蓝色作为打底
bg_img = Image.new('RGB', (width, height), (8, 9, 15)) 

# 创建一个透明图层用来画发光的光球
glow_img = Image.new('RGBA', (width, height), (0, 0, 0, 0))
draw = ImageDraw.Draw(glow_img)

# 画几个不同颜色的大圆，模拟 Apple 经典的弥散光/极光效果
draw.ellipse((-400, -400, 800, 800), fill=(110, 40, 210, 140)) # 左上角：深紫罗兰
draw.ellipse((1100, 400, 2400, 1600), fill=(255, 30, 120, 110)) # 右下角：洋红/深粉
draw.ellipse((300, 700, 1200, 1500), fill=(0, 150, 255, 100))   # 底部中央：湖蓝

# 施加强烈的高斯模糊，让色块融合成高级的流光渐变
glow_img = glow_img.filter(ImageFilter.GaussianBlur(200))

# 合并图层并保存
bg_img.paste(glow_img, (0,0), glow_img)
bg_img.save("gorgeous_bg.png")

print("2. 背景渲染完成，开始排版华丽的演示文稿...")
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
blank_layout = prs.slide_layouts[6]

# 定义主题颜色：香槟金、纯白、浅灰
GOLD = RGBColor(212, 175, 55)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(220, 220, 230)

def add_gorgeous_slide(title, subtitle="", slide_type="cover"):
    slide = prs.slides.add_slide(blank_layout)
    
    # 插入刚刚生成的流光背景
    slide.shapes.add_picture("gorgeous_bg.png", 0, 0, width=Inches(16), height=Inches(9))
    
    if slide_type == "cover":
        # 封面页：增加一个优雅的金色极细线框增加奢华感
        frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.5), Inches(13), Inches(6))
        frame.fill.background() # 内部透明
        frame.line.color.rgb = GOLD
        frame.line.width = Pt(1.5)
        
        # 主标题
        txBox = slide.shapes.add_textbox(Inches(2), Inches(3.2), Inches(12), Inches(2))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        # 衬线字体增加高级感
        p.font.name = 'Georgia' 
        p.font.size = Pt(85)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        # 副标题
        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.alignment = PP_ALIGN.CENTER
            p2.font.name = 'Helvetica Neue'
            p2.font.size = Pt(36)
            p2.font.color.rgb = GOLD
            p2.space_before = Pt(30)
            
    elif slide_type == "content":
        # 内容页排版
        txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(13), Inches(1.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.LEFT
        p.font.name = 'Georgia'
        p.font.size = Pt(60)
        p.font.bold = True
        p.font.color.rgb = GOLD
        
        # 标题下方的金色分割短线
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.8), Inches(2.5), Pt(3))
        line.fill.solid()
        line.fill.fore_color.rgb = GOLD
        line.line.fill.background()
        
        if subtitle:
            txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.6), Inches(13), Inches(4))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = subtitle
            p2.alignment = PP_ALIGN.LEFT
            p2.font.name = 'Helvetica Neue'
            p2.font.size = Pt(44)
            p2.font.color.rgb = WHITE
            p2.line_spacing = 1.5
            
    elif slide_type == "quote":
        # 金句页，居中排版，斜体
        txBox = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(12), Inches(4))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f'"{title}"'
        p.alignment = PP_ALIGN.CENTER
        p.font.name = 'Georgia'
        p.font.size = Pt(55)
        p.font.italic = True
        p.font.color.rgb = WHITE
        
        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.alignment = PP_ALIGN.RIGHT
            p2.font.name = 'Helvetica Neue'
            p2.font.size = Pt(32)
            p2.font.color.rgb = GOLD
            p2.space_before = Pt(40)

# 生成幻灯片内容
add_gorgeous_slide("公共关系实务", "建 立 信 任 的 高 级 艺 术", "cover")
add_gorgeous_slide("重塑认知", "公关不仅是危机的“灭火器”，更是品牌价值的“扩音器”。\n我们管理的不是信息，而是公众的认知。", "content")
add_gorgeous_slide("品牌金字塔", "从知名度、美誉度到忠诚度。\n你的品牌，就是你在别人心中的倒影。", "content")
add_gorgeous_slide("危机公关术", "黄金24小时已被打破，现在是“黄金1小时”。\n在风暴中心，真诚是唯一的避风港。", "content")
add_gorgeous_slide("媒介生态学", "传统媒体定调，社交媒体引爆。\n让真实的声音，以最优美的姿态触达受众。", "content")
add_gorgeous_slide("建立名誉需要20年，毁掉它只需5分钟。\n如果你明白这一点，你做事的态度就会截然不同。", "— 沃伦·巴菲特", "quote")
add_gorgeous_slide("致谢", "沟 通 · 信 任 · 共 赢", "cover")

output_file = 'PR_Practice_Gorgeous_Style.pptx'
prs.save(output_file)
print(f"3. 华丽版幻灯片生成成功: {output_file}")
