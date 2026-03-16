from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

print(">>> 启动《公共关系实务(第四版)》定制课件渲染引擎...")

def render_fluid_bg(filename, base_color, orbs):
    width, height = 1920, 1080
    bg = Image.new('RGB', (width, height), base_color)
    glow = Image.new('RGBA', (width, height), (0, 0, 0, 0))
    draw = ImageDraw.Draw(glow)
    for orb in orbs:
        x, y, r, color = orb
        draw.ellipse((x-r, y-r, x+r, y+r), fill=color)
    glow = glow.filter(ImageFilter.GaussianBlur(300))
    bg.paste(glow, (0,0), glow)
    bg.save(filename)

# 专属情绪配色
bgs = {
    "cover": ((8, 10, 15), [(960, 540, 800, (60, 100, 220, 140)), (400, 200, 800, (100, 200, 250, 100))]),
    "slide1": ((10, 10, 10), [(-200, 540, 900, (120, 120, 140, 100)), (1600, 100, 700, (80, 180, 200, 90))]),
    "slide2": ((5, 10, 15), [(1600, 540, 1000, (40, 160, 220, 120))]),
    "slide3": ((15, 8, 5), [(1600, 540, 1000, (255, 120, 50, 120))]),
    "slide4": ((15, 0, 5), [(1600, 540, 1000, (220, 40, 80, 140))]),
    "quote": ((0, 0, 0), [(960, 1080, 800, (200, 200, 220, 60))])
}

for name, data in bgs.items():
    render_fluid_bg(f"bg_{name}.png", data[0], data[1])

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
blank_layout = prs.slide_layouts[6]

WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(150, 150, 155)
DARK_GRAY = RGBColor(35, 35, 40)
ACCENT = RGBColor(100, 200, 255) # 科技蓝点缀

def add_text(slide, left, top, width, height, text, font_size, color, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = 'Helvetica Neue'
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return p

# --- 1. 封面页 ---
slide = prs.slides.add_slide(blank_layout)
slide.shapes.add_picture("bg_cover.png", 0, 0, width=Inches(16), height=Inches(9))
add_text(slide, Inches(1), Inches(3.2), Inches(14), Inches(2), "公共关系实务", 96, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.2), Inches(14), Inches(1), "（第四版）", 40, WHITE, False, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(6.5), Inches(14), Inches(1), "主编：霍瑞红   |   项目化导学课件", 28, ACCENT, False, PP_ALIGN.CENTER)

# --- 通用高级排版函数 ---
def add_course_slide(bg_name, number, title, subtitle, points):
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.add_picture(f"bg_{bg_name}.png", 0, 0, width=Inches(16), height=Inches(9))
    
    # 背景巨型浮水印数字
    add_text(slide, Inches(0.5), Inches(0.5), Inches(8), Inches(8), number, 400, DARK_GRAY, True)
    
    # 标题区
    add_text(slide, Inches(6.5), Inches(1.5), Inches(8.5), Inches(1.5), title, 60, WHITE, True)
    add_text(slide, Inches(6.5), Inches(2.8), Inches(8.5), Inches(1), subtitle, 32, ACCENT, False)
    
    # 知识点区
    txBox = slide.shapes.add_textbox(Inches(6.5), Inches(4.5), Inches(8.5), Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, pt in enumerate(points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = pt
        p.font.name = 'Helvetica Neue'
        p.font.size = Pt(26)
        p.font.color.rgb = WHITE if ("：" in pt or "项目" in pt) else GRAY
        p.space_after = Pt(20)
        p.line_spacing = 1.3

# --- 2. 教材特色与大纲 ---
add_course_slide("slide1", "00", "课程体系", "项目导向 · 任务驱动", [
    "本教材（霍瑞红主编）采用先进的项目化教学体系，包含：",
    "• 基础篇：项目一 (公关认知) / 项目二 (组织与公众)",
    "• 核心篇：项目三 (传播沟通) / 项目四 (调查与策划) / 项目五 (实施评估)",
    "• 实战篇：项目六 (专题活动) / 项目七 (危机公关) / 项目八 (公关礼仪)"
])

# --- 3. 项目一：公关认知 ---
add_course_slide("slide2", "01", "公关认知", "重新定义公共关系", [
    "• 经典定义：公共关系（PR）是社会组织为了塑造良好形象，通过传播沟通手段来影响公众的科学与艺术。",
    "• 核心本质：不仅是一项管理职能，更是组织与环境之间的“润滑剂”。",
    "• 职业素养：要求从业者具备敏锐的信息捕捉能力、卓越的表达能力与危机意识。"
])

# --- 4. 核心考点：三要素 ---
add_course_slide("slide3", "02", "基本要素", "公关运作的底层逻辑", [
    "公共关系活动由三大基本要素构成，缺一不可：",
    "1. 主体 (Subject)：社会组织（谁来做公关？）",
    "2. 客体 (Object)：公众（对谁做公关？）",
    "3. 手段 (Medium)：传播沟通（用什么方式做？）",
    "没有任何公关活动能够脱离这三个支柱独立存在。"
])

# --- 5. 项目七：高阶实战预告 ---
add_course_slide("slide4", "07", "危机公关", "全书实战难点与高潮", [
    "• 危机预防：在日常中建立预警机制，化解潜在风险。",
    "• 处理原则：真诚沟通、速度第一、承担责任、权威证实。",
    "• 实战任务：撰写危机公关声明、模拟召开新闻发布会应对媒体连番质询。"
])

# --- 6. 尾页 ---
slide = prs.slides.add_slide(blank_layout)
slide.shapes.add_picture("bg_quote.png", 0, 0, width=Inches(16), height=Inches(9))
add_text(slide, Inches(2), Inches(3.5), Inches(12), Inches(2), "理论指导实战\n沟通创造价值", 64, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(2), Inches(6.5), Inches(12), Inches(1), "《公共关系实务》（第四版）教学课件", 24, GRAY, False, PP_ALIGN.CENTER)

output_file = 'Huo_PR_Courseware_Pro.pptx'
prs.save(output_file)
print(f">>> 定制教材版课件渲染完成！文件已保存为: {output_file}")
