from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

print(">>> 启动 Ultra 级幻灯片渲染引擎...")

# 生成极致平滑的流体渐变背景图 (每页不同)
def render_fluid_bg(filename, base_color, orbs):
    width, height = 1920, 1080
    bg = Image.new('RGB', (width, height), base_color)
    glow = Image.new('RGBA', (width, height), (0, 0, 0, 0))
    draw = ImageDraw.Draw(glow)
    
    for orb in orbs:
        x, y, r, color = orb
        draw.ellipse((x-r, y-r, x+r, y+r), fill=color)
        
    # 使用极高的高斯模糊，制造极致柔和的 Apple 级渐变网格
    glow = glow.filter(ImageFilter.GaussianBlur(300))
    bg.paste(glow, (0,0), glow)
    bg.save(filename)

# 定义各页面的情绪配色方案 (x, y, radius, (R, G, B, Alpha))
bgs = {
    "cover": ((5, 5, 10), [ # 深空蓝紫
        (960, 540, 700, (88, 86, 214, 180)), 
        (400, 200, 800, (255, 45, 85, 120)),
        (1500, 800, 600, (90, 200, 250, 100))
    ]),
    "slide1": ((8, 8, 8), [ # 银河灰/青
        (-200, 540, 900, (100, 100, 120, 150)),
        (1600, 100, 700, (50, 200, 200, 90))
    ]),
    "slide2": ((10, 5, 0), [ # 璀璨金/橙 (品牌)
        (1600, 540, 1000, (255, 149, 0, 140)),
        (1800, 900, 600, (255, 204, 0, 100))
    ]),
    "slide3": ((15, 0, 5), [ # 警示红/洋红 (危机)
        (1600, 540, 1000, (255, 59, 48, 160)),
        (1000, 1080, 800, (255, 45, 85, 120))
    ]),
    "slide4": ((0, 10, 10), [ # 极光绿/青 (媒介)
        (1600, 540, 1000, (52, 199, 89, 130)),
        (1200, 0, 800, (90, 200, 250, 110))
    ]),
    "quote": ((0, 0, 0), [ # 纯黑深邃
        (960, 1080, 800, (200, 200, 220, 80))
    ])
}

# 渲染所有背景
print("1. 正在计算光线追踪级别的流体渐变纹理...")
for name, data in bgs.items():
    render_fluid_bg(f"bg_{name}.png", data[0], data[1])

print("2. 正在合成超高精度排版...")
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
blank_layout = prs.slide_layouts[6]

WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(142, 142, 147) # Apple System Gray
DARK_GRAY = RGBColor(40, 40, 45) # 水印数字颜色

def add_text(slide, left, top, width, height, text, font_size, color, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = 'Helvetica Neue'
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return p

# --- 封面 ---
slide = prs.slides.add_slide(blank_layout)
slide.shapes.add_picture("bg_cover.png", 0, 0, width=Inches(16), height=Inches(9))
add_text(slide, Inches(1), Inches(3), Inches(14), Inches(2), "公共关系实务", 96, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5), Inches(14), Inches(1), "Public Relations: The Art of Trust", 32, GRAY, False, PP_ALIGN.CENTER)

# --- 巨型数字分屏排版生成器 ---
def add_split_slide(bg_name, number, title, subtitle, desc):
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.add_picture(f"bg_{bg_name}.png", 0, 0, width=Inches(16), height=Inches(9))
    
    # 背景巨型浮水印数字 (极具视觉冲击力)
    add_text(slide, Inches(0.5), Inches(0.5), Inches(8), Inches(8), number, 400, DARK_GRAY, True)
    
    # 右侧内容区 (非对称高级排版)
    add_text(slide, Inches(7), Inches(2.5), Inches(8), Inches(1.5), title, 72, WHITE, True)
    add_text(slide, Inches(7), Inches(4.2), Inches(8), Inches(1), subtitle, 36, WHITE, False)
    
    # 正文描述
    p = add_text(slide, Inches(7), Inches(5.5), Inches(8), Inches(3), desc, 24, GRAY, False)
    p.line_spacing = 1.5

# 生成核心内容页
add_split_slide("slide1", "01", "重塑认知", "不仅是发通稿，更是管理公众认知。", 
                "公关的本质是沟通与连接。在信息爆炸的时代，\n我们不是在创造事实，而是在管理公众对事实的感知。\n这需要极度的敏锐与同理心。")

add_split_slide("slide2", "02", "品牌共鸣", "你的品牌，是你在别人心中的倒影。", 
                "从知名度、美誉度到忠诚度，这是一场漫长的攀登。\n形象塑造不仅是视觉传达，更是价值观的持续输出与验证。")

add_split_slide("slide3", "03", "风暴之眼", "危机公关：在极速崩塌中寻找生机。", 
                "黄金24小时已被打破，如今是“黄金1小时”。\n在风暴中心，唯有极度的真诚与果断的行动，\n才是度过危机的唯一避风港。")

add_split_slide("slide4", "04", "媒介生态", "让真实的声音，以最优美的姿态触达。", 
                "传统媒体定调，社交媒体引爆。\n理解媒介的底层逻辑，才能在嘈杂的互联网生态中，\n建立起属于你的信任灯塔。")

# --- 金句/尾页 ---
slide = prs.slides.add_slide(blank_layout)
slide.shapes.add_picture("bg_quote.png", 0, 0, width=Inches(16), height=Inches(9))
p = add_text(slide, Inches(2), Inches(3), Inches(12), Inches(3), "“建立名誉需要20年，毁掉它只需5分钟。”", 54, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(2), Inches(4.8), Inches(12), Inches(1), "— 沃伦·巴菲特", 28, GRAY, False, PP_ALIGN.CENTER)

output_file = 'PR_Practice_Ultra_Pro.pptx'
prs.save(output_file)
print(f">>> 终极渲染完成！文件已保存为: {output_file}")
