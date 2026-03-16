import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
# Set aspect ratio to 16:9 for modern displays (typical Apple Keynote size)
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Layout 6 is usually "Blank", giving us full control
blank_layout = prs.slide_layouts[6]

def add_apple_slide(title_text, subtitle_text=""):
    slide = prs.slides.add_slide(blank_layout)
    
    # Set black background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0) # True Black

    # Add text box in the center
    left = Inches(1)
    top = Inches(3.5) if not subtitle_text else Inches(3)
    width = Inches(14)
    height = Inches(2)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    # Main Title
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    # Use Helvetica or PingFang (Standard clean sans-serif fonts)
    p.font.name = 'Helvetica Neue'
    # Huge size for impact
    p.font.size = Pt(80) if len(title_text) < 15 else Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255) # White text
    
    # Subtitle (Optional)
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.alignment = PP_ALIGN.CENTER
        p2.font.name = 'Helvetica Neue'
        p2.font.size = Pt(40)
        p2.font.bold = False
        p2.font.color.rgb = RGBColor(160, 160, 160) # Apple Gray text
        p2.space_before = Pt(30)
        
    return slide

# Generating the Presentation Slides
add_apple_slide("公共关系实务", "建立信任的艺术")
add_apple_slide("不仅是发通稿", "更是认知管理")
add_apple_slide("核心一：形象塑造", "你的品牌，你的名片")
add_apple_slide("核心二：危机公关", "在风暴中寻找生机")
add_apple_slide("核心三：媒介沟通", "传递真实的声音")
add_apple_slide("“建立名誉需要20年，", "毁掉它只需5分钟。”\n— 沃伦·巴菲特")
add_apple_slide("沟通 · 信任 · 共赢", "Thank You")

prs.save('PR_Practice_Apple_Style.pptx')
print("Successfully generated PR_Practice_Apple_Style.pptx")
