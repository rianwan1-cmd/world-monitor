import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Set 16:9 ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Theme Colors
    ACADEMIC_BLUE = RGBColor(24, 50, 87)       # Deep Oxford Blue
    GOLDEN_YELLOW = RGBColor(212, 175, 55)     # Elegant Gold
    DARK_TEXT = RGBColor(50, 50, 50)
    CASE_RED = RGBColor(192, 0, 0)             # Engaging Red for cases

    def add_academic_background(slide):
        # Header bar
        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = ACADEMIC_BLUE
        header_shape.line.fill.background()

        # Golden accent line
        line_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, Inches(1.2), prs.slide_width, Inches(0.08)
        )
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = GOLDEN_YELLOW
        line_shape.line.fill.background()

        # Footer
        footer_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.4), prs.slide_width, Inches(0.4)
        )
        footer_shape.fill.solid()
        footer_shape.fill.fore_color.rgb = ACADEMIC_BLUE
        footer_shape.line.fill.background()
        
        # Footer text
        footer_box = slide.shapes.add_textbox(Inches(0.5), prs.slide_height - Inches(0.4), Inches(12), Inches(0.4))
        tf = footer_box.text_frame
        p = tf.paragraphs[0]
        p.text = "《公共管理学实务》（第四版）主编：霍瑞红 —— 理论与案例结合教学体系"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = "SimHei"

    # --- Title Slide ---
    slide_layout = prs.slide_layouts[6] # Blank
    title_slide = prs.slides.add_slide(slide_layout)
    
    # Background for title slide
    bg_shape = title_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = ACADEMIC_BLUE
    bg_shape.line.fill.background()
    
    # Title decorations
    dec_line = title_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.2), Inches(0.1), Inches(3))
    dec_line.fill.solid()
    dec_line.fill.fore_color.rgb = GOLDEN_YELLOW
    dec_line.line.fill.background()

    # Title Text
    title_box = title_slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(10), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "《公共管理学实务》第四版"
    p.alignment = PP_ALIGN.LEFT
    p.font.bold = True
    p.font.size = Pt(56)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = "SimHei"
    
    p_sub = tf.add_paragraph()
    p_sub.text = "17课时精讲（案例演练版）"
    p_sub.alignment = PP_ALIGN.LEFT
    p_sub.font.bold = True
    p_sub.font.size = Pt(44)
    p_sub.font.color.rgb = GOLDEN_YELLOW
    p_sub.font.name = "SimHei"

    subtitle_box = title_slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(10), Inches(1))
    tf2 = subtitle_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "主编：霍瑞红 | 风格：学院风·活泼生动"
    p2.alignment = PP_ALIGN.LEFT
    p2.font.size = Pt(28)
    p2.font.color.rgb = RGBColor(200, 200, 200)
    p2.font.name = "KaiTi"

    # --- Lessons Content ---
    lessons = [
        {"title": "第1课：公共管理导论与转型", "theory": "从传统行政到现代公共管理，强调效率、公平与服务导向，由“划桨”转向“掌舵”。", "case": "某市“最多跑一次”改革实务：如何打通政务服务最后一公里，用数据跑路代替群众跑腿？", "question": "思考互动：改革背后的部门利益阻力主要来自哪里？怎么破局？"},
        {"title": "第2课：公共管理者的角色与技能", "theory": "现代公共管理者不仅是政策执行者，更是利益协调者、创新者与危机应对者。", "case": "“接诉即办”机制下的社区干部实战：面对漏水纠纷，如何进行跨部门协调与资源调度？", "question": "情景模拟：如果你是社区书记，面对城管、住建相互推诿，你该怎么办？"},
        {"title": "第3课：公共政策的制定与执行", "theory": "政策过程包含问题界定、议程设置、方案规划、合法化、执行与评估六大环节。", "case": "共享单车管理政策的演变：从最初的放任不管，到一刀切清退，再到如今的电子围栏精细化治理。", "question": "小组讨论：政策执行中为什么容易出现“一刀切”的懒政现象？"},
        {"title": "第4课：公共组织设计与变革", "theory": "科层制（官僚制）的效率优势与僵化困境，扁平化、网络化、弹性化组织的兴起。", "case": "“大部制”改革的得与失：XX市市场监督管理局的整合之路，物理合并易，化学反应难。", "question": "头脑风暴：机构合并等于功能整合吗？组织文化如何融合？"},
        {"title": "第5课：公共人力资源管理实务", "theory": "从传统人事管理走向现代人力资源开发，强调战略匹配、激励机制与科学绩效评估。", "case": "聘任制公务员的探索：高达百万年薪引进高层次金融、大数据专业人才，鲶鱼效应还是水土不服？", "question": "辩论赛：公共部门的“铁饭碗”思维到底该不该彻底打破？"},
        {"title": "第6课：公共财政与预算管理", "theory": "公共财政的阳光化与法治化，预算是公共政策的数字化体现，反映政府的真实偏好。", "case": "参与式预算的社区实践：100万社区微更新资金，由居民投票决定是建车棚还是建花园。", "question": "现场投票：普通公众是否有能力参与复杂的政府预算编制流程？"},
        {"title": "第7课：信息资源管理与电子政务", "theory": "数字化转型重塑政府业务流程，打破信息孤岛，提升政府透明度与政民互动性。", "case": "智慧城市“大脑”建设：XX市“一网通办”背后的数据共享博弈与安全防御体系构建。", "question": "热点思考：政务数据共享面临的最大安全与隐私保护挑战是什么？"},
        {"title": "第8课：公共服务提供与社会化", "theory": "公共服务供给侧的多元化，引入市场机制与社会力量，政府购买服务成为常态。", "case": "养老服务的公私合作（PPP）模式：XX市社区食堂从排长队到面临生存危机的破局之路。", "question": "深度剖析：如何平衡公共服务的“公益性”与参与企业的“逐利性”？"},
        {"title": "第9课：非政府组织（NGO）与协同治理", "theory": "第三部门的崛起，填补政府失灵与市场失灵的空白，形成多元共治的现代治理格局。", "case": "环保NGO在跨流域治理中的作用：民间环保组织如何通过公益诉讼推动地方政府治污。", "question": "角色扮演：NGO如何既做政府的“牛虻”监督者，又做政府的合作伙伴？"},
        {"title": "第10课：危机管理与突发事件应对", "theory": "危机生命周期理论（4R模型）：缩减（Reduction）、预备（Readiness）、反应（Response）、恢复（Recovery）。", "case": "河南特大暴雨应急响应复盘反思：红色预警信息为何未能及时转化为有效的停工避险行动？", "question": "实战推演：突发事件初期的信息公开“黄金4小时”，首场发布会怎么开？"},
        {"title": "第11课：公共部门绩效评估实务", "theory": "以结果和民众满意度为导向，平衡计分卡（BSC）与关键绩效指标（KPI）在公共部门的应用。", "case": "某区政府KPI考核体系的优化升级：从过去单纯“看台账、看材料”转变为“看民意、看现场”。", "question": "反思质疑：过度依赖量化考核，会不会反而导致公共服务的形式主义？"},
        {"title": "第12课：公共管理伦理与责任", "theory": "公共利益至上原则，公职人员的职业道德规范，防范利益冲突与行政腐败的制度建设。", "case": "“旋转门”现象剖析：离职官员进入曾监管的利益相关企业高管层，合法合规与伦理争议的冲突。", "question": "灵魂拷问：防范腐败，严密的制度设计与内心的道德信念哪个更可靠？"},
        {"title": "第13课：城市治理与社区网格化管理", "theory": "城市即网格，社区即单元。网格化管理的精细化与城市治理的温度如何平衡。", "case": "老旧小区改造中的居民自治难题：“既有住宅加装电梯”中一楼与高层住户的利益协调。", "question": "基层调研：基层网格员是打通神经末梢的“全能战士”还是被形式主义困扰的“表哥表姐”？"},
        {"title": "第14课：乡村振兴与基层公共治理", "theory": "构建自治、法治、德治“三治结合”的现代乡村治理体系，实现乡村全面振兴。", "case": "某村“积分制”管理实践：用小积分撬动乡村大治理，好婆媳、垃圾分类都能换积分超市商品。", "question": "长效机制探讨：如果财政补贴的物质奖励取消，积分制还能长效运转吗？"},
        {"title": "第15课：全球化背景下的公共管理变革", "theory": "全球治理机制的构建，跨国界公共问题（如气候变化、传染病）对传统民族国家管理权限的挑战。", "case": "跨国公共卫生危机应对体系：以COVID-19为例，看国际组织协调与各国出入境管理政策的博弈。", "question": "宏观视野：主权国家在面对全球危机时，应如何有条件地让渡部分管理权以换取合作？"},
        {"title": "第16课：数字化与人工智能时代的治理创新", "theory": "算法治理、大数据决策的应用前景，以及随之而来的算法偏见、数字鸿沟与科技伦理挑战。", "case": "“算法官僚主义”：外卖骑手劳动权益保护与平台推荐算法的监管盲区，政府如何穿透算法？", "question": "前沿辩论：AI系统如果在公共决策（如福利发放审核）中出错，谁来承担法律责任？"},
        {"title": "第17课：课程总复习与实务大演练", "theory": "理论联系实际，知行合一。公共管理不是象牙塔里的学问，而是一门需要高度行动力与同理心的实践学科。", "case": "全班实战演练——大型模拟听证会：关于XX市是否应当在核心城区全面实施“拥堵费”的听证与决策。", "question": "结课总结：作为未来的公共管理者，这17堂课赋予了你哪些最核心的实务破局能力？"},
    ]

    for index, lesson in enumerate(lessons):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_academic_background(slide)

        # Content Box Background Shape
        content_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.8))
        content_bg.fill.solid()
        content_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        content_bg.line.fill.solid()
        content_bg.line.color.rgb = ACADEMIC_BLUE
        content_bg.line.width = Pt(1.5)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = lesson['title']
        p.font.size = Pt(38)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = "SimHei"
        p.alignment = PP_ALIGN.CENTER

        # Theory Section
        theory_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(11), Inches(1.2))
        tf_theory = theory_box.text_frame
        tf_theory.word_wrap = True
        p_theory = tf_theory.paragraphs[0]
        p_theory.text = "📚 理论精要：" 
        p_theory.font.size = Pt(26)
        p_theory.font.bold = True
        p_theory.font.color.rgb = ACADEMIC_BLUE
        p_theory.font.name = "SimHei"
        
        p_theory_text = tf_theory.add_paragraph()
        p_theory_text.text = lesson['theory']
        p_theory_text.font.size = Pt(22)
        p_theory_text.font.color.rgb = DARK_TEXT
        p_theory_text.font.name = "SimHei"

        # Case Section (Highlighting the "生动案例")
        case_box = slide.shapes.add_textbox(Inches(1.2), Inches(3.4), Inches(11), Inches(1.5))
        tf_case = case_box.text_frame
        tf_case.word_wrap = True
        p_case = tf_case.paragraphs[0]
        p_case.text = "🔥 实务案例实战："
        p_case.font.size = Pt(28)
        p_case.font.bold = True
        p_case.font.color.rgb = CASE_RED 
        p_case.font.name = "SimHei"
        
        p_case_text = tf_case.add_paragraph()
        p_case_text.text = lesson['case']
        p_case_text.font.size = Pt(26)
        p_case_text.font.bold = True
        p_case_text.font.color.rgb = DARK_TEXT
        p_case_text.font.name = "SimHei"

        # Question Section (Interactive part)
        q_box = slide.shapes.add_textbox(Inches(1.2), Inches(5.0), Inches(11), Inches(1))
        tf_q = q_box.text_frame
        tf_q.word_wrap = True
        p_q = tf_q.paragraphs[0]
        p_q.text = "💡 " + lesson['question']
        p_q.font.size = Pt(24)
        p_q.font.color.rgb = RGBColor(0, 100, 0) # Dark green for action
        p_q.font.name = "KaiTi"

    prs.save("Public_Management_Practice_Huo.pptx")
    print("PPT generated successfully!")

if __name__ == "__main__":
    create_presentation()
